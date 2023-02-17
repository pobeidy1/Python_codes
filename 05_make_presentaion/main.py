#import slide_data
from pptx import Presentation
from pptx.util import Inches

# slide_data = {
#     "slide1": {
#         "title": "Slide 1 Title",
#         "subtitle": "Slide 1 Subtitle",
#         "subtitle_bullet_list": ["Item 1", "Item 2", "Item 3"],
#         "figure": "slide1.png"
#     },
#     "slide2": {
#         "title": "Slide 2 Title",
#         "subtitle": "Slide 2 Subtitle",
#         "subtitle_bullet_list": ["Item A", "Item B", "Item C"],
#         "figure": "slide2.png"
#     }
# }

slide_data = {
    "slide1": {
        "title": "Introduction to Medical Image Formats",
        "subtitle": "Overview of Medical Image Formats",
        "subtitle_bullet_list": ["Purpose of Medical Image Formats", "Standardization of Medical Image Formats", "DICOM and NIfTI: Two Key Medical Image Formats"],
        "figure": "slide1.png"
    },
    "slide2": {
        "title": "Understanding DICOM Format",
        "subtitle": "",
        "subtitle_bullet_list": ["Definition of DICOM Format", "DICOM File Structure", "DICOM Header Information", "DICOM in Medical Imaging Workflow"],
        "figure": "slide2.png"
    },
    "slide3": {
        "title": "Understanding NIfTI Format",
        "subtitle": "",
        "subtitle_bullet_list": ["Definition of NIfTI Format", "NIfTI File Structure", "NIfTI Meta-Information", "Uses of NIfTI in Medical Imaging"],
        "figure": "slide3.png"
    },
    "slide4": {
        "title": "Converting DICOM to NIfTI",
        "subtitle": "Reasons for Conversion",
        "subtitle_bullet_list": ["Advantages of NIfTI Format for Image Processing", "Data De-identification", "File Size Reduction"],
        "figure": "slide4.png"
    },
    "slide5": {
        "title": "Converting DICOM to NIfTI",
        "subtitle": "Process Overview",
        "subtitle_bullet_list": ["Extracting Meta-Information from DICOM Header", "Resampling of Image Data", "Storing Image Data in NIfTI Format"],
        "figure": "slide5.png"
    },
    "slide6": {
        "title": "Orientation Information in NIfTI Format",
        "subtitle": "",
        "subtitle_bullet_list": ["Conventional Coordinate System", "Eigenvectors", "Orientation Information for Each Image"],
        "figure": "slide6.png"
    },
    "slide7": {
        "title": "Example: Stacking Multiple Images in NIfTI Format",
        "subtitle": "",
        "subtitle_bullet_list": ["Extracting Orientation Information for Each Image", "Concatenating Image Data", "Storing Result in NIfTI Format"],
        "figure": "slide7.png"
    },
    "slide8": {
        "title": "DICOM to NIfTI Conversion in Python",
        "subtitle": "",
        "subtitle_bullet_list": ["Introduction to dicom2nifti Package", "Reading DICOM Files", "Converting DICOM to NIfTI", "Saving NIfTI File"],
        "figure": "slide8.png"
    },
    "slide9": {
        "title": "Applications of NIfTI Format in AI",
       

    }}
# print(dir(slide_data))
# print(slide_data)
# slides = list(slide_data.items())
# # Create a new presentation object
# prs = Presentation()

import pptx

# Open the presentation
prs = pptx.Presentation()

# Loop through the slides
for slide_num, slide_info in slide_data():
    # Create a new slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Add the slide title
    title = slide.shapes.title
    title.text = slide_info["title"]

    # Add the slide subtitle
    subtitle = slide.placeholders[1]
    subtitle.text = slide_info["subtitle"]

    # Add the bullet points to the slide
    bullet_slide_layout = slide.slide_layout
    bullet_shape = bullet_slide_layout.shapes.placeholders[1].text_frame
    tf = bullet_shape.text_frame
    for point in slide_info["bullet_list"]:
        tf.add_paragraph().text = point

    try:
        # Add the slide figure
        picture = slide.shapes.add_picture(slide_info["figure"], left, top, height, width)
    except FileNotFoundError:
        print("Error: figure file not found for slide {}".format(slide_num))

# Save the presentation
prs.save("presentation.pptx")
